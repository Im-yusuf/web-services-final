"""
Generates the EcoNest COMP3011 presentation as a .pptx file.
Run: python3 docs/generate_presentation.py
Output: docs/EcoNest_Presentation.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Colour palette ────────────────────────────────────────────────────────────
BG       = RGBColor(0x0F, 0x17, 0x2A)   # dark navy
CARD     = RGBColor(0x1A, 0x23, 0x3A)   # card background
ACCENT   = RGBColor(0x38, 0xBD, 0xF8)   # sky-blue accent
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
MUTED    = RGBColor(0x94, 0xA3, 0xB8)
GOLD     = RGBColor(0xFB, 0xBF, 0x24)
GREEN    = RGBColor(0x34, 0xD3, 0x99)
RED      = RGBColor(0xF8, 0x71, 0x71)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs: Presentation):
    """Add a fully blank layout slide."""
    blank = prs.slide_layouts[6]
    return prs.slides.add_slide(blank)


def fill_bg(slide, colour: RGBColor = BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = colour


def add_rect(slide, l, t, w, h, fill: RGBColor, alpha=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)   # MSO_SHAPE_TYPE.RECTANGLE
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    return shape


def add_text(slide, text, l, t, w, h,
             size=18, bold=False, colour=WHITE,
             align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = colour
    return tb


def h_line(slide, y, colour=ACCENT, width=Pt(1.5)):
    from pptx.util import Pt as Pt2
    line = slide.shapes.add_shape(1, Inches(0.4), y, Inches(12.53), Pt(2))
    line.fill.solid()
    line.fill.fore_color.rgb = colour
    line.line.fill.background()
    return line


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ═══════════════════════════════════════════════════════════════════════════════

def slide_title(prs):
    """Slide 1 – Title"""
    sl = blank_slide(prs)
    fill_bg(sl)
    # Gradient-ish accent bar left side
    add_rect(sl, 0, 0, Inches(0.08), SLIDE_H, ACCENT)
    # Big title
    add_text(sl, "EcoNest", Inches(0.5), Inches(1.6), Inches(8), Inches(1.5),
             size=60, bold=True, colour=WHITE)
    add_text(sl, "UK Housing Market Analytics API",
             Inches(0.5), Inches(2.9), Inches(10), Inches(0.8),
             size=26, bold=False, colour=ACCENT)
    add_text(sl, "COMP3011 — Web Services & Web Data  |  March 2026",
             Inches(0.5), Inches(3.7), Inches(10), Inches(0.5),
             size=16, colour=MUTED)
    # Tag pills
    for i, tag in enumerate(["Express + TypeScript", "PostgreSQL + Prisma",
                              "Vue 3 + Pinia", "Land Registry Data"]):
        add_rect(sl, Inches(0.5 + i * 2.7), Inches(4.6), Inches(2.5), Inches(0.45), CARD)
        add_text(sl, tag, Inches(0.55 + i * 2.7), Inches(4.62), Inches(2.4), Inches(0.4),
                 size=13, colour=ACCENT, align=PP_ALIGN.CENTER)
    # Bottom note
    add_text(sl, "HM Land Registry Price Paid Data  •  2024–2025  •  Full-stack deployed application",
             Inches(0.5), Inches(6.6), Inches(12), Inches(0.5),
             size=11, colour=MUTED, align=PP_ALIGN.CENTER)


def slide_overview(prs):
    """Slide 2 – Project Overview"""
    sl = blank_slide(prs)
    fill_bg(sl)
    add_rect(sl, 0, 0, Inches(0.08), SLIDE_H, ACCENT)
    add_text(sl, "Project Overview", Inches(0.5), Inches(0.2), Inches(10), Inches(0.6),
             size=32, bold=True, colour=WHITE)
    h_line(sl, Inches(0.95))

    # Left column – what is it
    add_text(sl, "What is EcoNest?", Inches(0.5), Inches(1.15), Inches(5.8), Inches(0.5),
             size=18, bold=True, colour=ACCENT)
    desc = (
        "A RESTful API and interactive web application that transforms raw "
        "HM Land Registry transaction records into actionable housing market "
        "intelligence.\n\n"
        "Users can explore price trends, regional heatmaps, affordability "
        "metrics, and save property listings — all backed by real 2024–2025 "
        "transaction data."
    )
    add_text(sl, desc, Inches(0.5), Inches(1.7), Inches(5.8), Inches(2.8),
             size=14, colour=WHITE)

    # Right column – why it matters
    add_rect(sl, Inches(6.8), Inches(1.1), Inches(6.1), Inches(5.5), CARD)
    add_text(sl, "Why This Matters", Inches(7.0), Inches(1.3), Inches(5.6), Inches(0.5),
             size=18, bold=True, colour=GOLD)
    bullets = [
        ("🏠", "Raw Land Registry CSVs are inaccessible to most users"),
        ("📊", "No free tool provides programmatic access to trend analytics"),
        ("💷", "Affordability is a critical policy and personal finance question"),
        ("🔍", "Multi-year data enables genuine growth rate calculations"),
        ("🛡️", "Authenticated multi-user system with real session management"),
    ]
    for i, (icon, text) in enumerate(bullets):
        add_text(sl, f"{icon}  {text}", Inches(7.0), Inches(1.9 + i * 0.85),
                 Inches(5.5), Inches(0.75), size=13, colour=WHITE)

    # Data source badge
    add_rect(sl, Inches(0.5), Inches(4.8), Inches(5.8), Inches(0.7), CARD)
    add_text(sl, "📁  Data Source: HM Land Registry Price Paid  •  ~500k+ transactions  •  2024–2025",
             Inches(0.6), Inches(4.85), Inches(5.6), Inches(0.55),
             size=12, colour=GREEN)


def slide_architecture(prs):
    """Slide 3 – Architecture"""
    sl = blank_slide(prs)
    fill_bg(sl)
    add_rect(sl, 0, 0, Inches(0.08), SLIDE_H, ACCENT)
    add_text(sl, "System Architecture", Inches(0.5), Inches(0.2), Inches(10), Inches(0.6),
             size=32, bold=True, colour=WHITE)
    h_line(sl, Inches(0.95))

    # Architecture boxes – left to right flow
    boxes = [
        ("👤  Browser", "Vue 3 SPA\nVercel CDN", ACCENT),
        ("⚡  API Layer", "Express + TypeScript\nRailway Cloud", GREEN),
        ("🗄️  Data Layer", "PostgreSQL 17\nPrisma ORM", GOLD),
    ]
    for i, (title, sub, col) in enumerate(boxes):
        x = Inches(0.6 + i * 4.2)
        add_rect(sl, x, Inches(1.4), Inches(3.6), Inches(1.6), CARD)
        add_text(sl, title, x + Inches(0.1), Inches(1.5), Inches(3.4), Inches(0.55),
                 size=16, bold=True, colour=col, align=PP_ALIGN.CENTER)
        add_text(sl, sub, x + Inches(0.1), Inches(2.0), Inches(3.4), Inches(0.8),
                 size=13, colour=WHITE, align=PP_ALIGN.CENTER)
        if i < 2:
            add_text(sl, "──────►", Inches(4.2 + i * 4.2), Inches(2.0), Inches(0.9), Inches(0.5),
                     size=20, colour=MUTED, align=PP_ALIGN.CENTER)

    # Layered breakdown
    add_text(sl, "Backend Request Pipeline", Inches(0.5), Inches(3.2), Inches(12), Inches(0.45),
             size=17, bold=True, colour=ACCENT)
    layers = [
        ("Routes", "URL pattern → controller mapping"),
        ("Middleware", "Auth check → Zod validation → error handler"),
        ("Controllers", "HTTP adapter: parse inputs, call service, format response"),
        ("Services", "Business logic: DB queries, aggregations, signal classification"),
    ]
    for i, (name, desc) in enumerate(layers):
        x = Inches(0.5 + i * 3.15)
        add_rect(sl, x, Inches(3.75), Inches(3.0), Inches(1.2), CARD)
        add_text(sl, name, x + Inches(0.1), Inches(3.82), Inches(2.8), Inches(0.4),
                 size=14, bold=True, colour=ACCENT, align=PP_ALIGN.CENTER)
        add_text(sl, desc, x + Inches(0.1), Inches(4.2), Inches(2.8), Inches(0.65),
                 size=11, colour=MUTED, align=PP_ALIGN.CENTER)
        if i < 3:
            add_text(sl, "→", Inches(3.5 + i * 3.15), Inches(4.1), Inches(0.3), Inches(0.4),
                     size=20, colour=MUTED, align=PP_ALIGN.CENTER)

    # Design decisions
    add_text(sl, "Key Design Decision: Controller–Service split keeps HTTP concerns out of business logic, enabling isolated testing.",
             Inches(0.5), Inches(5.2), Inches(12.3), Inches(0.5),
             size=12, colour=MUTED, align=PP_ALIGN.CENTER)

    add_rect(sl, Inches(0.5), Inches(5.85), Inches(12.3), Inches(0.75), CARD)
    add_text(sl, "Session tokens stored in DB (not JWTs) → instant revocation on logout  •  CORS: whitelist + *.vercel.app regex for preview deploys  •  Global error handler prevents stack trace leakage",
             Inches(0.6), Inches(5.9), Inches(12.1), Inches(0.65),
             size=11, colour=GREEN, align=PP_ALIGN.CENTER)


def slide_stack(prs):
    """Slide 4 – Technology Stack"""
    sl = blank_slide(prs)
    fill_bg(sl)
    add_rect(sl, 0, 0, Inches(0.08), SLIDE_H, ACCENT)
    add_text(sl, "Technology Stack & Justification", Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
             size=32, bold=True, colour=WHITE)
    h_line(sl, Inches(0.95))

    stack = [
        ("Node.js + Express", "TypeScript", "Async I/O for DB-heavy workloads; full end-to-end type safety; shared types between backend and frontend; npm ecosystem depth"),
        ("Prisma ORM", "Schema-first", "Type-safe queries generated from schema; migration history for reproducible deploys; clean groupBy/aggregate API without raw SQL risks"),
        ("PostgreSQL 17", "Relational DB", "Composite indexes on 7 columns power sub-second analytics queries; ACID + unique constraints for idempotent imports; advanced aggregation"),
        ("Vue 3 + Pinia", "Composition API", "Reactive state across 5 stores (auth, cache, filters, saved, comparison); dataCache store prevents redundant API calls"),
        ("Zod Validation", "Runtime safety", "Schema validation at API boundaries only — no over-engineering of internal code paths; clear error messages to clients"),
        ("Vitest + Supertest", "Integration tests", "Full HTTP layer tested against real DB; chosen over unit tests because query construction bugs only surface at integration level"),
    ]
    for i, (tech, tag, reason) in enumerate(stack):
        row, col = divmod(i, 2)
        x = Inches(0.5 + col * 6.4)
        y = Inches(1.15 + row * 1.75)
        add_rect(sl, x, y, Inches(6.1), Inches(1.55), CARD)
        add_text(sl, tech, x + Inches(0.15), y + Inches(0.1), Inches(3.5), Inches(0.4),
                 size=15, bold=True, colour=ACCENT)
        add_rect(sl, x + Inches(3.8), y + Inches(0.08), Inches(2.1), Inches(0.35), BG)
        add_text(sl, tag, x + Inches(3.82), y + Inches(0.1), Inches(2.0), Inches(0.32),
                 size=11, colour=GOLD, align=PP_ALIGN.CENTER)
        add_text(sl, reason, x + Inches(0.15), y + Inches(0.55), Inches(5.8), Inches(0.9),
                 size=11, colour=MUTED)

    add_text(sl, "Alternatives evaluated: Python FastAPI, Go Fiber, React, Sequelize, JWT auth — rejected with documented rationale in technical report.",
             Inches(0.5), Inches(6.95), Inches(12.3), Inches(0.4),
             size=11, colour=MUTED, align=PP_ALIGN.CENTER)


def slide_features(prs):
    """Slide 5 – Key Features"""
    sl = blank_slide(prs)
    fill_bg(sl)
    add_rect(sl, 0, 0, Inches(0.08), SLIDE_H, ACCENT)
    add_text(sl, "Key Features", Inches(0.5), Inches(0.2), Inches(10), Inches(0.6),
             size=32, bold=True, colour=WHITE)
    h_line(sl, Inches(0.95))

    features = [
        ("📈  Price Trends", "Monthly avg price + MoM growth %\nFilter by region, year, property type"),
        ("🗺️  Regional Heatmap", "County-level avg price + transaction\nvolume from DB groupBy aggregation"),
        ("🔍  Property Browser", "Paginated search: price range, type,\nregion — capped at 100/page"),
        ("💡  Market Insights", "9 parallel DB queries: affordability\nratio, 1yr/5yr growth, market signal"),
        ("🔖  Saved Listings", "Authenticated CRUD: save, annotate,\nand manage bookmarked properties"),
        ("⚖️  Comparison View", "Side-by-side compare up to N properties\nvia client-side Pinia store"),
        ("🤖  Market Assistant", "AI-powered summary of regional\nmarket conditions (Insights API)"),
        ("📊  Dashboard KPIs", "Stat cards: total properties, avg price,\nregions tracked, latest data date"),
    ]
    for i, (title, desc) in enumerate(features):
        row, col = divmod(i, 4)
        x = Inches(0.4 + col * 3.12)
        y = Inches(1.15 + row * 2.6)
        add_rect(sl, x, y, Inches(2.95), Inches(2.35), CARD)
        add_text(sl, title, x + Inches(0.12), y + Inches(0.12), Inches(2.7), Inches(0.55),
                 size=13, bold=True, colour=ACCENT)
        add_text(sl, desc, x + Inches(0.12), y + Inches(0.7), Inches(2.7), Inches(1.5),
                 size=11, colour=WHITE)


def slide_api(prs):
    """Slide 6 – API Design"""
    sl = blank_slide(prs)
    fill_bg(sl)
    add_rect(sl, 0, 0, Inches(0.08), SLIDE_H, ACCENT)
    add_text(sl, "API Design & Documentation", Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
             size=32, bold=True, colour=WHITE)
    h_line(sl, Inches(0.95))

    # Endpoint table
    add_text(sl, "Endpoints", Inches(0.5), Inches(1.1), Inches(7.5), Inches(0.4),
             size=18, bold=True, colour=ACCENT)
    endpoints = [
        ("POST /api/auth/register", "201", "No", "Register new user, returns session token"),
        ("POST /api/auth/login", "200", "No", "Authenticate, returns session token"),
        ("GET  /api/auth/me", "200", "✓ Yes", "Return authenticated user profile"),
        ("GET  /api/trends", "200", "No", "Monthly avg price + growth (filterable)"),
        ("GET  /api/heatmap", "200", "No", "County-level avg price + volume"),
        ("GET  /api/properties", "200", "No", "Paginated property listing with filters"),
        ("GET  /api/stats", "200", "No", "Total properties, avg price, latest date"),
        ("GET  /api/regions", "200", "No", "Distinct region names for filter UI"),
        ("POST /api/insights/region", "200", "No", "Full market profile for a region"),
        ("GET/POST/PUT/DELETE /api/saved", "200/201", "✓ Yes", "Saved listings CRUD"),
    ]
    headers = ["  Endpoint", "Status", "Auth", "Description"]
    col_widths = [Inches(3.1), Inches(0.7), Inches(0.7), Inches(4.4)]
    col_x = [Inches(0.4), Inches(3.55), Inches(4.3), Inches(5.05)]

    # Header row
    for j, (hdr, cx, cw) in enumerate(zip(headers, col_x, col_widths)):
        add_rect(sl, cx, Inches(1.55), cw - Inches(0.05), Inches(0.35), ACCENT)
        add_text(sl, hdr, cx + Inches(0.05), Inches(1.57), cw, Inches(0.32),
                 size=11, bold=True, colour=BG)

    for i, (ep, status, auth, desc) in enumerate(endpoints):
        y = Inches(1.95 + i * 0.38)
        bg = CARD if i % 2 == 0 else BG
        add_rect(sl, Inches(0.4), y, Inches(9.1), Inches(0.36), bg)
        cols = [ep, status, auth, desc]
        for j, (val, cx, cw) in enumerate(zip(cols, col_x, col_widths)):
            col = GREEN if auth == "✓ Yes" and j == 2 else (GOLD if j == 1 else WHITE)
            add_text(sl, val, cx + Inches(0.05), y + Inches(0.02), cw, Inches(0.32),
                     size=10, colour=col)

    # Right panel – response envelope + Swagger
    add_rect(sl, Inches(9.65), Inches(1.1), Inches(3.45), Inches(5.8), CARD)
    add_text(sl, "Response Envelope", Inches(9.8), Inches(1.2), Inches(3.2), Inches(0.4),
             size=14, bold=True, colour=ACCENT)
    code = ('{ "success": true,\n  "data": { ... } }\n\n'
            '{ "success": false,\n  "error": "message" }\n\n'
            'Auth: Bearer <token>\nin Authorization header\n\n'
            '7-day session expiry\nInstant revocation on logout')
    add_text(sl, code, Inches(9.8), Inches(1.65), Inches(3.15), Inches(3.8),
             size=11, colour=GREEN)
    add_text(sl, "📄  Swagger UI: /api/docs\nOpenAPI JSON: /api/docs.json",
             Inches(9.8), Inches(5.5), Inches(3.15), Inches(0.8),
             size=11, colour=GOLD)

    add_text(sl, "All inputs validated with Zod middleware  •  Consistent error codes follow RFC 7231  •  Parameterised queries prevent SQL injection",
             Inches(0.4), Inches(7.05), Inches(12.5), Inches(0.35),
             size=10, colour=MUTED, align=PP_ALIGN.CENTER)


def slide_data_model(prs):
    """Slide 7 – Data Model"""
    sl = blank_slide(prs)
    fill_bg(sl)
    add_rect(sl, 0, 0, Inches(0.08), SLIDE_H, ACCENT)
    add_text(sl, "Data Model & Schema", Inches(0.5), Inches(0.2), Inches(10), Inches(0.6),
             size=32, bold=True, colour=WHITE)
    h_line(sl, Inches(0.95))

    # Tables
    tables = [
        ("users", WHITE, ["id (UUID PK)", "email (unique)", "password_hash", "created_at"]),
        ("sessions", ACCENT, ["id (UUID PK)", "token (unique, indexed)", "user_id → users", "expires_at"]),
        ("property_sales", GOLD, ["transaction_id (unique)", "price, transfer_date", "postcode, town_city, district, county", "property_type, new_build, tenure", "7× indexes for query performance"]),
        ("saved_listings", GREEN, ["id (UUID PK)", "user_id → users", "property_id → property_sales", "notes (optional)", "UNIQUE(user_id, property_id)"]),
    ]
    positions = [(0.4, 1.1), (3.5, 1.1), (6.6, 1.1), (0.4, 3.9)]
    for (name, col, fields), (x, y) in zip(tables, positions):
        w = Inches(2.9) if name != "property_sales" else Inches(6.1)
        add_rect(sl, Inches(x), Inches(y), w, Inches(2.5 if name != "property_sales" else 2.5), CARD)
        add_rect(sl, Inches(x), Inches(y), w, Inches(0.4), col)
        add_text(sl, name, Inches(x + 0.1), Inches(y + 0.04), w - Inches(0.2), Inches(0.35),
                 size=14, bold=True, colour=BG, align=PP_ALIGN.CENTER)
        for fi, field in enumerate(fields):
            add_text(sl, f"  {field}", Inches(x + 0.1), Inches(y + 0.48 + fi * 0.38),
                     w - Inches(0.2), Inches(0.36), size=11, colour=WHITE)

    # Indexes highlight
    add_rect(sl, Inches(0.4), Inches(6.6), Inches(6.1), Inches(0.7), CARD)
    add_text(sl, "⚡  property_sales indexes: postcode • town_city • district • county • transfer_date • property_type • price → sub-second analytics queries",
             Inches(0.55), Inches(6.65), Inches(5.8), Inches(0.6), size=11, colour=GREEN)

    # Relationships
    add_rect(sl, Inches(6.8), Inches(3.9), Inches(6.1), Inches(3.35), CARD)
    add_text(sl, "Relationships & Integrity", Inches(6.95), Inches(4.0), Inches(5.8), Inches(0.45),
             size=14, bold=True, colour=ACCENT)
    rels = [
        "User ──1:*──► Sessions  (CASCADE DELETE)",
        "User ──1:*──► SavedListings  (CASCADE DELETE)",
        "PropertySale ──1:*──► SavedListings",
        "User ──1:1──► UserPreference",
        "UNIQUE (user_id, property_id) prevents duplicate saves",
        "UNIQUE transaction_id → idempotent data imports",
        "UUIDs as PKs throughout (no sequential ID leakage)",
    ]
    for i, r in enumerate(rels):
        add_text(sl, r, Inches(6.95), Inches(4.55 + i * 0.37), Inches(5.8), Inches(0.35),
                 size=11, colour=WHITE)


def slide_security(prs):
    """Slide 8 – Security"""
    sl = blank_slide(prs)
    fill_bg(sl)
    add_rect(sl, 0, 0, Inches(0.08), SLIDE_H, ACCENT)
    add_text(sl, "Security Implementation", Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
             size=32, bold=True, colour=WHITE)
    h_line(sl, Inches(0.95))

    measures = [
        ("🔑  Password Hashing", "bcrypt — 12 salt rounds\nPlaintext never stored or logged\nIndustry standard for credential storage"),
        ("🎲  Token Generation", "crypto.randomBytes(48) → 384-bit entropy\nCryptographically secure random tokens\nFar beyond brute-force feasibility"),
        ("⏱️  Session Expiry", "7-day TTL enforced server-side\nExpired sessions rejected even with valid token\nLogout immediately deletes DB row"),
        ("✅  Input Validation", "Zod schemas on all mutating endpoints\nValidation at system boundary only\nPrevents malformed data reaching DB"),
        ("🛡️  SQL Injection", "Prisma parameterised queries throughout\nNo raw SQL string construction\nORM type safety as second defence layer"),
        ("🕵️  User Enumeration", '"Invalid email or password" — same\nmessage for both failure modes\nPrevents attacker learning valid emails'),
        ("🌐  CORS Policy", "Explicit origin whitelist\n+ *.vercel.app regex for previews\nCredentials: true only for listed origins"),
        ("📋  Error Handling", "Global error handler catches everything\nStack traces never sent to clients\nConsistent { success: false, error } shape"),
    ]
    for i, (title, desc) in enumerate(measures):
        row, col = divmod(i, 4)
        x = Inches(0.4 + col * 3.12)
        y = Inches(1.15 + row * 2.55)
        add_rect(sl, x, y, Inches(2.95), Inches(2.3), CARD)
        add_text(sl, title, x + Inches(0.12), y + Inches(0.12), Inches(2.7), Inches(0.55),
                 size=12, bold=True, colour=GREEN)
        add_text(sl, desc, x + Inches(0.12), y + Inches(0.7), Inches(2.7), Inches(1.5),
                 size=11, colour=WHITE)

    add_text(sl, "OWASP Top 10 reviewed: no broken access control, no injection, no cryptographic failures, no server-side request forgery",
             Inches(0.4), Inches(7.05), Inches(12.5), Inches(0.35),
             size=10, colour=MUTED, align=PP_ALIGN.CENTER)


def slide_testing(prs):
    """Slide 9 – Testing"""
    sl = blank_slide(prs)
    fill_bg(sl)
    add_rect(sl, 0, 0, Inches(0.08), SLIDE_H, ACCENT)
    add_text(sl, "Testing Strategy", Inches(0.5), Inches(0.2), Inches(10), Inches(0.6),
             size=32, bold=True, colour=WHITE)
    h_line(sl, Inches(0.95))

    # Philosophy box
    add_rect(sl, Inches(0.5), Inches(1.1), Inches(12.3), Inches(0.85), CARD)
    add_text(sl, "💡  Philosophy: Integration tests over unit tests — query construction bugs only surface when the full HTTP → Service → DB pipeline runs together. Mocking the database hides the most common bug class in a data API.",
             Inches(0.65), Inches(1.15), Inches(12.0), Inches(0.75), size=12, colour=GREEN)

    # Test suites
    suites = [
        ("auth.test.ts", ACCENT, [
            "✓  Register new user (201 + token)",
            "✓  Reject duplicate email (409)",
            "✓  Reject invalid email / short password (400)",
            "✓  Login with valid credentials (200 + token)",
            "✓  Reject wrong password / unknown user (401)",
            "✓  GET /me with valid token (200)",
            "✓  Reject /me without token (401)",
            "✓  Logout invalidates session",
        ]),
        ("properties.test.ts", GOLD, [
            "✓  Trends endpoint — no filters",
            "✓  Trends — filtered by region + year + type",
            "✓  Heatmap — data shape validation",
            "✓  Properties — paginated default",
            "✓  Properties — with price/type filters",
            "✓  Stats — total, avg, latest date",
            "✓  Regions — returns array",
        ]),
        ("savedListings.test.ts", GREEN, [
            "✓  Auth guard on all saved routes",
            "✓  Create saved listing (201)",
            "✓  List saved listings (200)",
            "✓  Update note on saved listing (200)",
            "✓  Delete saved listing (200)",
            "✓  Reject duplicate save (409)",
        ]),
    ]
    for i, (name, col, cases) in enumerate(suites):
        x = Inches(0.4 + i * 4.2)
        add_rect(sl, x, Inches(2.1), Inches(4.0), Inches(4.85), CARD)
        add_rect(sl, x, Inches(2.1), Inches(4.0), Inches(0.38), col)
        add_text(sl, name, x + Inches(0.1), Inches(2.12), Inches(3.8), Inches(0.35),
                 size=13, bold=True, colour=BG, align=PP_ALIGN.CENTER)
        for j, case in enumerate(cases):
            add_text(sl, case, x + Inches(0.12), Inches(2.58 + j * 0.38),
                     Inches(3.75), Inches(0.36), size=10, colour=WHITE)

    add_text(sl, "Stack: Vitest + Supertest  •  npm test from backend/  •  TypeScript type-check: npx tsc --noEmit",
             Inches(0.4), Inches(7.05), Inches(12.5), Inches(0.35),
             size=10, colour=MUTED, align=PP_ALIGN.CENTER)


def slide_deployment(prs):
    """Slide 10 – Deployment & Struggles"""
    sl = blank_slide(prs)
    fill_bg(sl)
    add_rect(sl, 0, 0, Inches(0.08), SLIDE_H, ACCENT)
    add_text(sl, "Deployment & Operations", Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
             size=32, bold=True, colour=WHITE)
    h_line(sl, Inches(0.95))

    # Deployment boxes
    deploys = [
        ("🔵  Railway", "Backend API", ["Express — auto-deploy from GitHub main", "PORT, DATABASE_URL, CORS_ORIGIN env vars", "railway.json: build + start + healthcheck", "Prisma migrations run on every deploy"]),
        ("🐘  Railway Postgres", "Database", ["PostgreSQL 17 managed service", "Internal URL for runtime (low latency)", "External URL for local data imports", "500k+ records imported via CLI script"]),
        ("⚡  Vercel", "Frontend SPA", ["Vue 3 + Vite — auto-deploy from GitHub", "VITE_API_BASE_URL → Railway backend", "vercel.json → SPA fallback routing", "Preview deployments per PR/commit"]),
    ]
    for i, (title, sub, bullets) in enumerate(deploys):
        x = Inches(0.4 + i * 4.3)
        add_rect(sl, x, Inches(1.1), Inches(4.1), Inches(3.0), CARD)
        add_text(sl, title, x + Inches(0.15), Inches(1.2), Inches(3.8), Inches(0.45),
                 size=15, bold=True, colour=ACCENT)
        add_text(sl, sub, x + Inches(0.15), Inches(1.65), Inches(3.8), Inches(0.35),
                 size=12, colour=MUTED)
        for j, b in enumerate(bullets):
            add_text(sl, f"→  {b}", x + Inches(0.15), Inches(2.1 + j * 0.4),
                     Inches(3.8), Inches(0.38), size=11, colour=WHITE)

    # Deployment struggles section
    add_text(sl, "⚠️  Deployment Struggles (Most Frustrating Part)", Inches(0.5), Inches(4.25), Inches(12), Inches(0.45),
             size=16, bold=True, colour=RED)
    struggles = [
        ("Railway build failures", "First deploys failed — Nixpacks couldn't resolve TS build. Had to iterate on railway.json config multiple times."),
        ("Internal vs External DB URLs", "Tried using Railway internal URL from my laptop for imports. Didn't work — had to learn the difference the hard way."),
        ("CORS preflight rejections", "Frontend on Vercel, backend on Railway = different domains. Browser blocked everything until regex origin check was added."),
        ("Slow remote data imports", "Importing 500k+ records over external DB URL was slow and timed out. Had to tune batch sizes carefully."),
    ]
    for i, (title, desc) in enumerate(struggles):
        row, col = divmod(i, 2)
        x = Inches(0.4 + col * 6.35)
        y = Inches(4.8 + row * 1.2)
        add_rect(sl, x, y, Inches(6.15), Inches(1.0), CARD)
        add_text(sl, title, x + Inches(0.1), y + Inches(0.05), Inches(5.9), Inches(0.35),
                 size=12, bold=True, colour=RED)
        add_text(sl, desc, x + Inches(0.1), y + Inches(0.4), Inches(5.9), Inches(0.55),
                 size=10, colour=MUTED)


def slide_genai(prs):
    """Slide 11 – GenAI Usage"""
    sl = blank_slide(prs)
    fill_bg(sl)
    add_rect(sl, 0, 0, Inches(0.08), SLIDE_H, ACCENT)
    add_text(sl, "Generative AI Usage — 3-Tool Pipeline", Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.6),
             size=28, bold=True, colour=WHITE)
    h_line(sl, Inches(0.95))

    # Pipeline banner
    add_rect(sl, Inches(0.5), Inches(1.1), Inches(12.3), Inches(0.75), CARD)
    add_text(sl, "[Me]  →  Gemini  (idea + requirements)  →  ChatGPT  (prompt engineering)  →  Copilot / Claude  (implementation)",
             Inches(0.65), Inches(1.18), Inches(12.0), Inches(0.58), size=14, bold=True, colour=GOLD, align=PP_ALIGN.CENTER)

    # Three pipeline boxes — more personal narrative
    pipeline_boxes = [
        ("1. Google Gemini", ACCENT,
         "I didn't know what to build.\n~15 min brainstorming converged\non housing market API.\n\nThen asked Gemini to produce\na full requirements document\ngiven my chosen tech stack."),
        ("2. ChatGPT", GREEN,
         'Gave it Gemini\'s requirements\nand asked: "Write me the best\nprompt for Claude to build this."\n\nChatGPT structured the prompt\nwith file layout, types, layering\nrules, security specs, tests.'),
        ("3. Copilot / Claude", GOLD,
         "Built the full project from\nthe ChatGPT-engineered prompt.\n\nI reviewed all code, debated\narchitecture (JWT → sessions),\nran tests, debugged deployment\nissues manually."),
    ]
    for i, (title, col, desc) in enumerate(pipeline_boxes):
        x = Inches(0.4 + i * 4.2)
        add_rect(sl, x, Inches(2.0), Inches(4.0), Inches(2.9), CARD)
        add_rect(sl, x, Inches(2.0), Inches(4.0), Inches(0.4), col)
        add_text(sl, title, x + Inches(0.1), Inches(2.02), Inches(3.8), Inches(0.37),
                 size=13, bold=True, colour=BG, align=PP_ALIGN.CENTER)
        add_text(sl, desc, x + Inches(0.12), Inches(2.5), Inches(3.75), Inches(2.3),
                 size=11, colour=WHITE)
        if i < 2:
            add_text(sl, "→", Inches(4.35 + i * 4.2), Inches(3.3), Inches(0.3), Inches(0.4),
                     size=22, colour=MUTED, align=PP_ALIGN.CENTER)

    # Honest reflection box
    add_rect(sl, Inches(0.5), Inches(5.05), Inches(12.3), Inches(0.9), CARD)
    add_text(sl, "Honest Reflection", Inches(0.65), Inches(5.1), Inches(3.0), Inches(0.35),
             size=13, bold=True, colour=RED)
    add_text(sl, "Most code was AI-generated from a prompt that was itself AI-written, from requirements that were AI-produced. My role: project director — chose the project, decided the stack, reviewed all code, made architectural decisions (sessions > JWTs), debugged deployment manually, ran tests, wrote analytical content. I understand every part of the system and can explain any component.",
             Inches(0.65), Inches(5.45), Inches(11.95), Inches(0.48),
             size=10, colour=MUTED)

    # Usage summary table
    usages = [
        ("Gemini", "Project ideation (~15 min chat)", "Chose housing market API over alternatives"),
        ("Gemini", "Requirements analysis", "Full functional + non-functional spec given my stack"),
        ("ChatGPT", "Prompt engineering", "Converted requirements → precise Copilot build prompt"),
        ("Copilot", "Full backend + frontend implementation", "Built from engineered prompt with my ongoing review"),
        ("Copilot", "Architecture refinement", "JWT→sessions debate; median vs mean; REST vs GraphQL"),
        ("Copilot", "Deployment debugging", "CORS regex fix; Railway config; DB URL confusion"),
        ("Copilot", "Test review + report/slides", "Identified edge cases; structured documentation"),
    ]
    add_rect(sl, Inches(0.5), Inches(6.05), Inches(12.3), Inches(0.3), ACCENT)
    for j, (hdr, cx, cw) in enumerate(zip(
        ["  Tool", "  Usage", "  Detail"],
        [Inches(0.55), Inches(1.85), Inches(6.6)],
        [Inches(1.25), Inches(4.7), Inches(6.0)]
    )):
        add_text(sl, hdr, cx, Inches(6.07), cw, Inches(0.26), size=10, bold=True, colour=BG)
    for i, (tool, use, detail) in enumerate(usages):
        y = Inches(6.39 + i * 0.22)
        bg = CARD if i % 2 == 0 else BG
        add_rect(sl, Inches(0.5), y, Inches(12.3), Inches(0.2), bg)
        tool_col = {"Gemini": ACCENT, "ChatGPT": GREEN, "Copilot": GOLD}.get(tool, WHITE)
        add_text(sl, f"  {tool}", Inches(0.55), y, Inches(1.25), Inches(0.2), size=8, colour=tool_col, bold=True)
        add_text(sl, f"  {use}", Inches(1.85), y, Inches(4.7), Inches(0.2), size=8, colour=WHITE)
        add_text(sl, f"  {detail}", Inches(6.6), y, Inches(6.0), Inches(0.2), size=8, colour=MUTED)


def slide_version_control(prs):
    """Slide 12 – Version Control"""
    sl = blank_slide(prs)
    fill_bg(sl)
    add_rect(sl, 0, 0, Inches(0.08), SLIDE_H, ACCENT)
    add_text(sl, "Version Control Practices", Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
             size=32, bold=True, colour=WHITE)
    h_line(sl, Inches(0.95))

    add_text(sl, "GitHub Repository", Inches(0.5), Inches(1.1), Inches(6.0), Inches(0.45),
             size=18, bold=True, colour=ACCENT)
    add_rect(sl, Inches(0.5), Inches(1.6), Inches(6.0), Inches(4.5), CARD)
    practices = [
        "✓  Public repository — visible commit history",
        "✓  Consistent commits throughout development",
        "✓  README.md with full setup instructions",
        "✓  .gitignore excludes node_modules, .env files",
        "✓  Separate backend/ and frontend/ workspaces",
        "✓  Prisma migrations tracked in version control",
        "✓  No secrets committed — all via env vars",
        "✓  Data files excluded (large CSVs not in repo)",
        "✓  Code corresponds to presented version",
    ]
    for i, p in enumerate(practices):
        col = GREEN if p.startswith("✓") else RED
        add_text(sl, p, Inches(0.65), Inches(1.7 + i * 0.45), Inches(5.6), Inches(0.42),
                 size=12, colour=col)

    add_text(sl, "Commit History Highlights", Inches(6.9), Inches(1.1), Inches(5.9), Inches(0.45),
             size=18, bold=True, colour=ACCENT)
    add_rect(sl, Inches(6.9), Inches(1.6), Inches(5.9), Inches(4.5), CARD)
    commits = [
        ("init", "Initial project structure, Express setup"),
        ("feat", "Prisma schema + migration for property_sales"),
        ("feat", "Auth service: bcrypt + crypto session tokens"),
        ("feat", "Property trends + heatmap endpoints"),
        ("feat", "Market insights: 9 parallel aggregations"),
        ("feat", "Saved listings CRUD with auth guard"),
        ("feat", "Vue 3 frontend: router, stores, views"),
        ("feat", "Comparison view + dataCache Pinia store"),
        ("test", "Vitest + Supertest integration test suites"),
        ("docs", "Swagger JSDoc on all route files"),
        ("deploy", "Railway + Vercel production config"),
    ]
    for i, (tag, msg) in enumerate(commits):
        y = Inches(1.7 + i * 0.37)
        add_rect(sl, Inches(7.05), y, Inches(0.7), Inches(0.3),
                 {"init": MUTED, "feat": ACCENT, "test": GREEN, "docs": GOLD, "deploy": RED}.get(tag, MUTED))
        add_text(sl, tag, Inches(7.05), y, Inches(0.7), Inches(0.3), size=9, bold=True,
                 colour=BG, align=PP_ALIGN.CENTER)
        add_text(sl, msg, Inches(7.8), y, Inches(4.8), Inches(0.32), size=10, colour=WHITE)

    add_rect(sl, Inches(0.5), Inches(6.25), Inches(12.3), Inches(0.6), CARD)
    add_text(sl, "All deliverables accessible from repo: README.md links to API docs, technical report, and presentation slides.",
             Inches(0.65), Inches(6.3), Inches(12.0), Inches(0.5), size=12, colour=GOLD, align=PP_ALIGN.CENTER)


def slide_challenges(prs):
    """Slide 13 – Challenges & Limitations"""
    sl = blank_slide(prs)
    fill_bg(sl)
    add_rect(sl, 0, 0, Inches(0.08), SLIDE_H, ACCENT)
    add_text(sl, "Challenges, Limitations & Future Work", Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.6),
             size=30, bold=True, colour=WHITE)
    h_line(sl, Inches(0.95))

    # Challenges
    add_text(sl, "Challenges Overcome", Inches(0.5), Inches(1.1), Inches(6.0), Inches(0.45),
             size=18, bold=True, colour=RED)
    challenges = [
        ("Deployment Pain", "Railway build failures, CORS rejections,\nDB URL confusion — most time-consuming\npart of the entire project."),
        ("Data Volume", "500k+ records — needed 7 indexes on\nfilter columns + batch import tuning\nfor acceptable query performance."),
        ("CORS Split-Deploy", "Vercel generates unique preview URLs.\nRegex origin check was the solution\nbut took hours to diagnose."),
        ("Session vs JWT", "JWTs can't be revoked on logout.\nSwitched to DB sessions — added\ncomplexity but real security benefit."),
        ("TypeScript Strictness", "Strict mode caught real bugs:\noptional params used as strings,\nnullable aggregate results."),
        ("AI Limitations", "AI was least helpful for deployment.\nCouldn't see Railway logs or Vercel\nbuild output — manual trial and error."),
    ]
    for i, (title, desc) in enumerate(challenges):
        row, col = divmod(i, 2)
        x = Inches(0.5 + col * 3.0)
        y = Inches(1.65 + row * 1.55)
        add_rect(sl, x, y, Inches(2.8), Inches(1.35), CARD)
        add_text(sl, title, x + Inches(0.1), y + Inches(0.08), Inches(2.6), Inches(0.35),
                 size=12, bold=True, colour=RED)
        add_text(sl, desc, x + Inches(0.1), y + Inches(0.45), Inches(2.6), Inches(0.85),
                 size=10, colour=WHITE)

    # Limitations + Future
    add_text(sl, "Limitations & Future Work", Inches(6.8), Inches(1.1), Inches(6.0), Inches(0.45),
             size=18, bold=True, colour=GOLD)
    limitations = [
        "⚠️  Median price is approximate (no SQL PERCENTILE)",
        "⚠️  No rate limiting on public endpoints yet",
        "⚠️  MarketReport cache model defined but unused",
        "⚠️  Heatmap is county-level text (no coordinates)",
        "⚠️  Static dataset — no incremental monthly updates",
    ]
    for i, lim in enumerate(limitations):
        add_text(sl, lim, Inches(6.85), Inches(1.65 + i * 0.45),
                 Inches(6.1), Inches(0.42), size=11, colour=GOLD)

    add_text(sl, "Future Development", Inches(6.8), Inches(3.95), Inches(6.0), Inches(0.45),
             size=18, bold=True, colour=GREEN)
    futures = [
        "🚀  PostGIS → real map-based heatmaps",
        "🚀  Cron job → monthly HMLR data updates",
        "🚀  Redis/materialised views → analytics caching",
        "🚀  express-rate-limit → production safety",
        "🚀  WebSocket → real-time price alerts",
    ]
    for i, f in enumerate(futures):
        add_text(sl, f, Inches(6.85), Inches(4.45 + i * 0.45),
                 Inches(6.1), Inches(0.42), size=11, colour=GREEN)


def slide_deliverables(prs):
    """Slide 14 – Deliverables Summary"""
    sl = blank_slide(prs)
    fill_bg(sl)
    add_rect(sl, 0, 0, Inches(0.08), SLIDE_H, ACCENT)
    add_text(sl, "Deliverables Summary", Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
             size=32, bold=True, colour=WHITE)
    h_line(sl, Inches(0.95))

    deliverables = [
        ("✅  Code Repository", GREEN, [
            "Public GitHub repository",
            "Visible commit history throughout development",
            "README.md with full setup instructions",
            "All code runnable — matches presented version",
            "backend/ and frontend/ workspaces",
        ]),
        ("✅  API Documentation", GREEN, [
            "Swagger UI at /api/docs (interactive)",
            "OpenAPI JSON at /api/docs.json",
            "Every endpoint documented with JSDoc",
            "Example requests and responses",
            "Auth process and error codes documented",
        ]),
        ("✅  Technical Report", GREEN, [
            "5-page report (docs/TECHNICAL_REPORT.md)",
            "Stack justification with alternatives considered",
            "Architecture, security, and testing sections",
            "Challenges, limitations, future work",
            "GenAI declaration with reflective analysis",
        ]),
        ("✅  Presentation Slides", GREEN, [
            "This presentation (PPTX)",
            "Version control section ✓",
            "API documentation section ✓",
            "Technical report highlights ✓",
            "All deliverables covered ✓",
        ]),
    ]
    for i, (title, col, items) in enumerate(deliverables):
        row, c = divmod(i, 2)
        x = Inches(0.5 + c * 6.4)
        y = Inches(1.1 + row * 2.85)
        add_rect(sl, x, y, Inches(6.1), Inches(2.65), CARD)
        add_text(sl, title, x + Inches(0.15), y + Inches(0.1), Inches(5.8), Inches(0.5),
                 size=16, bold=True, colour=col)
        for j, item in enumerate(items):
            add_text(sl, f"  {item}", x + Inches(0.15), y + Inches(0.65 + j * 0.38),
                     Inches(5.8), Inches(0.35), size=11, colour=WHITE)

    add_text(sl, "All deliverables linked from README.md  •  Report submitted as PDF  •  Slides submitted as PPTX  •  Repo is public",
             Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.35),
             size=11, colour=MUTED, align=PP_ALIGN.CENTER)


def slide_qa(prs):
    """Slide 15 – Q&A"""
    sl = blank_slide(prs)
    fill_bg(sl)
    add_rect(sl, 0, 0, Inches(0.08), SLIDE_H, ACCENT)
    add_text(sl, "Thank You", Inches(0.5), Inches(1.4), Inches(12), Inches(1.2),
             size=54, bold=True, colour=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, "Questions & Answers", Inches(0.5), Inches(2.9), Inches(12), Inches(0.7),
             size=28, colour=ACCENT, align=PP_ALIGN.CENTER)

    h_line(sl, Inches(3.8))

    # Anticipated Q&A
    add_text(sl, "Anticipated Questions", Inches(0.5), Inches(4.0), Inches(12), Inches(0.45),
             size=16, bold=True, colour=GOLD, align=PP_ALIGN.CENTER)

    qnas = [
        ("Why sessions over JWTs?",
         "Instant revocation on logout — JWTs can't be invalidated without a blocklist. One extra DB read per auth check is acceptable with indexed token column."),
        ("What was your biggest struggle?",
         "Deployment. Railway build failures, CORS preflight rejections with Vercel, DB URL confusion (internal vs external). AI couldn't help much here — manual trial and error."),
        ("How did you use 3 different AI tools?",
         "Gemini for ideation + requirements, ChatGPT to engineer a precise build prompt, Copilot/Claude to implement. Each tool chosen for its strength. I reviewed everything."),
        ("How do you handle 500k+ records?",
         "7 indexes on filter columns for fast queries. Batch import (500/batch) with streaming CSV parser. In-DB aggregation for heatmap, in-app for trend calculations."),
    ]
    for i, (q, a) in enumerate(qnas):
        row, col = divmod(i, 2)
        x = Inches(0.5 + col * 6.4)
        y = Inches(4.55 + row * 1.25)
        add_rect(sl, x, y, Inches(6.1), Inches(1.15), CARD)
        add_text(sl, f"Q: {q}", x + Inches(0.15), y + Inches(0.08), Inches(5.8), Inches(0.4),
                 size=12, bold=True, colour=ACCENT)
        add_text(sl, a, x + Inches(0.15), y + Inches(0.5), Inches(5.8), Inches(0.6),
                 size=10, colour=WHITE)


# ═══════════════════════════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════════════════════════

def build():
    prs = new_prs()

    slide_title(prs)
    slide_overview(prs)
    slide_architecture(prs)
    slide_stack(prs)
    slide_features(prs)
    slide_api(prs)
    slide_data_model(prs)
    slide_security(prs)
    slide_testing(prs)
    slide_deployment(prs)
    slide_genai(prs)
    slide_version_control(prs)
    slide_challenges(prs)
    slide_deliverables(prs)
    slide_qa(prs)

    out = os.path.join(os.path.dirname(__file__), "EcoNest_Presentation.pptx")
    prs.save(out)
    print(f"✅  Saved: {out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
